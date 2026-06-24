import os
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

from langchain_core.documents import Document

# --- Configuration & Paths ---
BASE_DIR = os.path.dirname(os.path.dirname(__file__))
TEXTBOOKS_DIR = os.path.join(BASE_DIR, "data", "textbooks")
VECTOR_STORE_DIR = os.path.join(BASE_DIR, "vector_store", "faiss_index")


def build_vector_database():
    """Loads PDFs, DOCX, PPTX, TXT, chunks the text, and builds a local FAISS database."""
    print("[*] Starting RAG Scholar Engine...")

    # 1. Load Documents
    if not os.path.exists(TEXTBOOKS_DIR):
        os.makedirs(TEXTBOOKS_DIR)
        print("[-] Error: 'data/textbooks/' folder is empty. Please add a document and run again.")
        return

    print(f"[*] Loading files from {TEXTBOOKS_DIR}...")
    documents = []
    
    for filename in os.listdir(TEXTBOOKS_DIR):
        file_path = os.path.join(TEXTBOOKS_DIR, filename)
        if not os.path.isfile(file_path):
            continue
            
        try:
            if filename.lower().endswith(".pdf"):
                import PyPDF2
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
                    if text.strip():
                        documents.append(Document(page_content=text, metadata={"source": filename}))
            elif filename.lower().endswith(".docx"):
                import docx
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                if text.strip():
                    documents.append(Document(page_content=text, metadata={"source": filename}))
            elif filename.lower().endswith(".pptx"):
                import pptx
                prs = pptx.Presentation(file_path)
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                if text.strip():
                    documents.append(Document(page_content=text, metadata={"source": filename}))
            elif filename.lower().endswith((".xlsx", ".xls")):
                try:
                    import pandas as pd
                    # Attempt to read all sheets and combine them
                    xls = pd.ExcelFile(file_path)
                    sheet_texts = []
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name=sheet_name)
                        sheet_texts.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
                    text = "\n\n".join(sheet_texts)
                    if text.strip():
                        documents.append(Document(page_content=text, metadata={"source": filename}))
                except Exception as e:
                    print(f"[-] Excel parsing failed for {filename}: {e}")
            elif filename.lower().endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
                    if text.strip():
                        documents.append(Document(page_content=text, metadata={"source": filename}))
        except Exception as e:
            print(f"[-] Failed to parse {filename}: {e}")

    if not documents:
        print("[-] No valid documents with text found in the directory. Add a file and try again.")
        return

    print(f"[*] Loaded {len(documents)} document(s).")

    # 2. Chunk the Text
    print("[*] Chunking text into readable segments...")
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )
    chunks = text_splitter.split_documents(documents)
    print(f"[*] Created {len(chunks)} text chunks.")

    if not chunks:
        print("[-] No chunks were created (documents might be empty or unreadable). Skipping FAISS index.")
        return

    # 3. Create Embeddings & Store in FAISS
    print("[*] Generating AI embeddings...")
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

    # Build the FAISS index
    vector_db = FAISS.from_documents(chunks, embeddings)

    # 4. Save Database Locally
    os.makedirs(VECTOR_STORE_DIR, exist_ok=True)
    vector_db.save_local(VECTOR_STORE_DIR)
    print(f"[*] Success! Vector database saved locally at {VECTOR_STORE_DIR}")


def search_textbooks(query, k=3):
    """Searches the local FAISS database for relevant paragraphs."""
    if not os.path.exists(VECTOR_STORE_DIR):
        return "Vector database not found. Please build it first."

    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

    # Load the saved database (Dangerous deserialization is disabled by default, we allow it for local use)
    vector_db = FAISS.load_local(VECTOR_STORE_DIR, embeddings, allow_dangerous_deserialization=True)

    # Perform semantic search
    results = vector_db.similarity_search(query, k=k)

    # Combine the top results into a single string
    context = "\n\n".join([doc.page_content for doc in results])
    return context


def search_textbooks_with_sources(query, k=3):
    """Searches the local FAISS database and returns both content and source list."""
    if not os.path.exists(VECTOR_STORE_DIR):
        return "Vector database not found. Please build it first.", []

    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    vector_db = FAISS.load_local(VECTOR_STORE_DIR, embeddings, allow_dangerous_deserialization=True)
    results = vector_db.similarity_search(query, k=k)

    context = "\n\n".join([doc.page_content for doc in results])
    sources = list(set([doc.metadata.get("source") for doc in results if doc.metadata.get("source")]))
    return context, sources


if __name__ == "__main__":
    # 1. Build the database (Run this whenever you add new PDFs)
    build_vector_database()

    # 2. Test the semantic search
    print("\n--- Testing Semantic Search ---")
    test_query = "What is the main topic of this document?"
    print(f"Query: {test_query}")
    print("Found Context:")
    print("--------------------------------------------------")
    context, sources = search_textbooks_with_sources(test_query, k=1)
    print(context)
    print("Sources:", sources)
    print("--------------------------------------------------")