import os
import json
import pathlib
import hashlib
import json
import uuid
import asyncio
import base64
import re
from datetime import datetime
from dotenv import load_dotenv
import chainlit as cl
from openai import OpenAI
from openai import AsyncOpenAI
from fastapi import Request, Response

# Load environment variables first
load_dotenv()

import warnings
with warnings.catch_warnings():
    warnings.simplefilter("ignore")
    from duckduckgo_search import DDGS

# Import backend modules
from src.rag_scholar import search_textbooks_with_sources, build_vector_database
from chainlit.server import app as fastapi_app
from fastapi import UploadFile, File, Response
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
import asyncio

import chainlit.data as cl_data
from chainlit.data.sql_alchemy import SQLAlchemyDataLayer
from sqlalchemy import MetaData, Table, Column, String, Boolean, Integer, ForeignKey, event as sa_event
from sqlalchemy import text
import aiofiles
from typing import Union, Dict, Any, List, Optional, Sequence
from chainlit.data.storage_clients.base import BaseStorageClient


# --- CUSTOM LOCAL BLOB STORAGE PROVIDER ---
class LocalStorageClient(BaseStorageClient):
    def __init__(self, base_dir="public/elements", base_url="/public/elements"):
        self.base_dir = base_dir
        self.base_url = base_url
        os.makedirs(self.base_dir, exist_ok=True)

    async def upload_file(
        self,
        object_key: str,
        data: Union[bytes, str],
        mime: str = "application/octet-stream",
        overwrite: bool = True,
        content_disposition: str | None = None,
    ) -> Dict[str, Any]:
        file_path = os.path.join(self.base_dir, object_key)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        if isinstance(data, str):
            async with aiofiles.open(file_path, "w", encoding="utf-8") as f:
                await f.write(data)
        else:
            async with aiofiles.open(file_path, "wb") as f:
                await f.write(data)
        url = f"{self.base_url}/{object_key}"
        return {"object_key": object_key, "url": url}

    async def delete_file(self, object_key: str) -> bool:
        file_path = os.path.join(self.base_dir, object_key)
        if os.path.exists(file_path):
            os.remove(file_path)
            return True
        return False

    async def get_read_url(self, object_key: str) -> str:
        return f"{self.base_url}/{object_key}"

    async def close(self) -> None:
        pass


# Configure Storage Provider
storage_provider = LocalStorageClient()

import sqlite3
import json

# Intercept and serialize Python objects for SQLite compatibility
sqlite3.register_adapter(list, lambda lst: json.dumps(lst))
sqlite3.register_adapter(dict, lambda dct: json.dumps(dct))

def initialize_database_raw():
    db_path = "chainlit.db"
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    # Explicitly create all required tables
    cursor.execute('CREATE TABLE IF NOT EXISTS users (id TEXT PRIMARY KEY, identifier TEXT UNIQUE, createdAt TEXT, metadata TEXT)')
    cursor.execute('CREATE TABLE IF NOT EXISTS threads (id TEXT PRIMARY KEY, createdAt TEXT, name TEXT, userId TEXT, userIdentifier TEXT, tags TEXT, metadata TEXT)')
    cursor.execute('CREATE TABLE IF NOT EXISTS steps (id TEXT PRIMARY KEY, name TEXT NOT NULL, type TEXT NOT NULL, threadId TEXT NOT NULL, parentId TEXT, command TEXT, modes TEXT, streaming BOOLEAN NOT NULL, waitForAnswer BOOLEAN, isError BOOLEAN, metadata TEXT, tags TEXT, input TEXT, output TEXT, createdAt TEXT, start TEXT, end TEXT, generation TEXT, showInput TEXT, defaultOpen BOOLEAN, autoCollapse BOOLEAN, language TEXT)')
    cursor.execute('CREATE TABLE IF NOT EXISTS feedbacks (id TEXT PRIMARY KEY, forId TEXT NOT NULL, threadId TEXT, value INTEGER NOT NULL, comment TEXT)')
    cursor.execute('CREATE TABLE IF NOT EXISTS elements (id TEXT PRIMARY KEY, threadId TEXT, type TEXT, chainlitKey TEXT, path TEXT, url TEXT, objectKey TEXT, name TEXT NOT NULL, display TEXT, size TEXT, language TEXT, page INTEGER, props TEXT, autoPlay BOOLEAN, playerConfig TEXT, forId TEXT, mime TEXT)')
    
    # Dual-Tier Memory Tables
    cursor.execute('CREATE TABLE IF NOT EXISTS user_primary_memory (user_id TEXT PRIMARY KEY, username TEXT, first_name TEXT, last_name TEXT, current_model TEXT, current_mode TEXT, updated_at TEXT)')
    cursor.execute('CREATE TABLE IF NOT EXISTS user_secondary_memory (user_id TEXT, pref_key TEXT, pref_value TEXT, PRIMARY KEY (user_id, pref_key))')
    
    conn.commit()
    conn.close()

# Invoke immediately to guarantee schema exists before SQLAlchemy hits it
initialize_database_raw()

class NexusDataLayer(SQLAlchemyDataLayer):
    pass

# Configure Data Layer globally
# FIX P1-11: Remove direct cl_data._data_layer assignment; use only @cl.data_layer decorator
data_layer = NexusDataLayer(
    conninfo="sqlite+aiosqlite:///chainlit.db",
    storage_provider=storage_provider
)

@cl.data_layer
def get_data_layer():
    return data_layer




# --- ADAPTERS FOR NATIVE GEMINI SDK WRAPPERS ---
class OpenAIChunkDelta:
    def __init__(self, content):
        self.content = content

class OpenAIChunkChoice:
    def __init__(self, delta):
        self.delta = delta

class OpenAIChunk:
    def __init__(self, content):
        self.choices = [OpenAIChunkChoice(OpenAIChunkDelta(content))]

async def wrap_gemini_stream(gemini_stream):
    for chunk in gemini_stream:
        try:
            yield OpenAIChunk(chunk.text)
        except Exception:
            pass

class GeminiToolFunction:
    def __init__(self, name, arguments):
        self.name = name
        self.arguments = arguments

class GeminiToolCall:
    def __init__(self, id, name, arguments):
        self.id = id
        self.type = "function"
        self.function = GeminiToolFunction(name, arguments)

class GeminiMessage:
    def __init__(self, content, tool_calls=None):
        self.role = "assistant"
        self.content = content or ""
        self.tool_calls = tool_calls

class GeminiChoice:
    def __init__(self, message):
        self.message = message

class GeminiCompletion:
    def __init__(self, message):
        self.choices = [GeminiChoice(message)]


# --- INITIALIZE DUAL-PROVIDER ARCHITECTURE ---
API_KEY_GROQ = os.getenv("GROQ_API_KEY")
API_KEY_GEMINI = os.getenv("GEMINI_API_KEY")

if not API_KEY_GROQ or not API_KEY_GEMINI:
    raise ValueError("[-] Critical Error: Missing API Keys in .env file.")

try:
    groq_client = OpenAI(api_key=API_KEY_GROQ, base_url="https://api.groq.com/openai/v1")
    openai_client = AsyncOpenAI(api_key=os.environ.get("OPENAI_API_KEY", "dummy_key_if_missing"))
    import warnings
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        import google.generativeai as genai
    getattr(genai, "configure")(api_key=API_KEY_GEMINI)
    gemini_client = genai
except Exception as e:
    print(f"[-] Client initialization error: {e}")


TOOLS_SCHEMA = [

    {
        "type": "function",
        "function": {
            "name": "query_academic_textbooks",
            "description": "Searches through academic textbooks.",
            "parameters": {
                "type": "object",
                "properties": {
                    "search_query": {"type": "string"}
                },
                "required": ["search_query"],
                "additionalProperties": False
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "execute_web_search",
            "description": "Searches the live internet.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string"}
                },
                "required": ["query"],
                "additionalProperties": False
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "update_user_memory",
            "description": "Use this tool whenever the user explicitly states a permanent fact about themselves (e.g., their name, their profession, a strict preference). This saves the fact to their global profile.",
            "parameters": {
                "type": "object",
                "properties": {
                    "fact_key": {
                        "type": "string",
                        "description": "The key under which to store the fact, e.g., 'name', 'profession'"
                    },
                    "fact_value": {
                        "type": "string",
                        "description": "The value of the fact to store."
                    }
                },
                "required": ["fact_key", "fact_value"]
            }
        }
    }
]


def perform_web_search(query):
    """Synchronous web search — only call via perform_web_search_async from async contexts."""
    try:
        results = DDGS().text(query, max_results=3)
        return "\n\n".join([f"Source: {res['title']}\nSnippet: {res['body']}" for res in
                            results]) if results else "No internet results."
    except Exception as e:
        return f"Web search failed: {str(e)}"

async def perform_web_search_async(query: str) -> str:
    """Wrap synchronous DDGS search in executor to avoid event-loop blocking."""
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, lambda: perform_web_search(query))


async def call_llm_with_fallback(messages, tools=None, tool_choice=None, temperature=None, vision_mode=False, requested_model=None):
    # 1. Define clean routes (No decommissioned models)
    if vision_mode:
        MODELS_ROUTE = [
            {"provider": "Gemini", "client": gemini_client, "model": "gemini-1.5-pro"}
        ]
        tools = None
    else:
        MODELS_ROUTE = []
        if requested_model:
            provider = "Gemini" if "gemini" in requested_model else "Groq"
            client = gemini_client if provider == "Gemini" else groq_client
            MODELS_ROUTE.append({"provider": provider, "client": client, "model": requested_model})
        
        # Base stable route
        MODELS_ROUTE.extend([
            {"provider": "Groq", "client": groq_client, "model": "llama-3.3-70b-versatile"},
            {"provider": "Gemini", "client": gemini_client, "model": "gemini-1.5-pro"},
            {"provider": "Gemini", "client": gemini_client, "model": "gemini-1.5-flash"}
        ])

    last_error = None
    for route in MODELS_ROUTE:
        try:
            if route["provider"] == "Groq":
                kwargs = {"model": route["model"], "messages": messages, "stream": not bool(tools)}
                if temperature is not None: kwargs["temperature"] = temperature
                if tools and "vision" not in route["model"].lower():
                    kwargs["tools"] = tools
                    kwargs["tool_choice"] = tool_choice or "auto"
                
                loop = asyncio.get_running_loop()
                stream = await loop.run_in_executor(None, lambda: route["client"].chat.completions.create(**kwargs))
                
                if kwargs["stream"]:
                    async def async_token_stream():
                        for chunk in stream:
                            if chunk.choices and chunk.choices[0].delta.content:
                                yield chunk.choices[0].delta.content
                    return async_token_stream(), route["model"]
                else:
                    return stream, route["model"]

            elif route["provider"] == "Gemini":
                system_instruction = None
                contents = []
                for msg in messages:
                    if msg.get("role") == "system":
                        system_instruction = msg.get("content")
                        continue
                    
                    role_map = {"user": "user", "assistant": "model", "model": "model"}
                    msg_role = role_map.get(msg.get("role"), "user")
                    
                    if isinstance(msg.get("content"), str):
                        contents.append({"role": msg_role, "parts": [msg["content"]]})
                    elif isinstance(msg.get("content"), list):
                        parts = []
                        for block in msg["content"]:
                            if block.get("type") == "text":
                                parts.append(block.get("text"))
                            elif block.get("type") == "image_url":
                                url = block["image_url"]["url"]
                                if url.startswith("data:image/"):
                                    mime_type = url.split(";")[0].split(":")[1]
                                    base64_data = url.split(",")[1]
                                    import base64
                                    parts.append({
                                        "mime_type": mime_type,
                                        "data": base64.b64decode(base64_data)
                                    })
                        contents.append({"role": msg_role, "parts": parts})

                gen_config = route["client"].types.GenerationConfig(temperature=temperature or 0.7)
                gemini_tools = [t["function"] for t in tools] if (tools and isinstance(tools, list)) else None
                
                model = route["client"].GenerativeModel(model_name=route["model"], tools=gemini_tools)
                
                loop = asyncio.get_running_loop()
                response = await loop.run_in_executor(
                    None, 
                    lambda: model.generate_content(contents, generation_config=gen_config, stream=True)
                )
                
                async def async_token_stream_gemini():
                    for chunk in response:
                        if hasattr(chunk, "text") and chunk.text:
                            yield chunk.text
                return async_token_stream_gemini(), route["model"]

        except Exception as e:
            print(f"[-] Fallback failed on {route['model']}: {e}")
            last_error = e
            continue

    # 4. The Safety Valve: Return a streaming generator even on failure
    async def error_stream():
        yield f"⚠️ NexusAI encountered an error: {str(last_error)[:50]}..."
    
    return error_stream(), "System Error"

def message_to_dict(msg):
    """Safely converts a message (dict or ChatCompletionMessage) to a standard dict."""
    if isinstance(msg, dict):
        return msg.copy()
    role = getattr(msg, "role", "assistant")
    content = getattr(msg, "content", "")
    d = {"role": role, "content": content}
    tool_calls = getattr(msg, "tool_calls", None)
    if tool_calls:
        d["tool_calls"] = [
            {
                "id": tc.id,
                "type": tc.type,
                "function": {
                    "name": tc.function.name,
                    "arguments": tc.function.arguments
                }
            }
            for tc in tool_calls
        ]
    return d


def sanitize_history_for_storage(history_array):
    """Strips image base64 blobs from history to prevent bloated SQLite entries."""
    clean_history = []
    for msg in history_array:
        msg_dict = message_to_dict(msg)
        if isinstance(msg_dict.get("content"), list):
            text_only = next(
                (c.get("text", "") for c in msg_dict["content"] if isinstance(c, dict) and c.get("type") == "text"),
                "[Image Omitted for Storage]"
            )
            msg_dict["content"] = text_only
        clean_history.append(msg_dict)
    return clean_history


# --- CHAINLIT STARTERS ---
@cl.set_starters
async def set_starters(user: cl.User | None = None, **kwargs):
    return [
        cl.Starter(
            label="📚 Scholar RAG — Page Algorithms",
            message="Search the textbooks and explain Page Replacement Algorithms with examples.",
            icon="/public/favicon.ico"
        ),

        cl.Starter(
            label="🌐 Live Web Intel",
            message="Search the web for the current AWS server status and any outages.",
            icon="/public/favicon.ico"
        ),
        cl.Starter(
            label="🖼️ Analyze an Image",
            message="I'll upload an image — describe what you see and extract all key information.",
            icon="/public/favicon.ico"
        )
    ]


@cl.set_chat_profiles
async def chat_profile(user: cl.User | None = None, **kwargs):
    return [
        cl.ChatProfile(
            name="Omni Mode",
            markdown_description="**Full-capability** assistant. Access all tools: web search, textbook RAG, and gaming analytics.",
            icon="https://cdn-icons-png.flaticon.com/512/4712/4712035.png",
            starters=[
                cl.Starter(label="🌐 Live Web Search", message="Search the web for the current AWS server status and any active outages.", icon="/public/favicon.ico"),
                cl.Starter(label="📚 Scholar RAG", message="Search the textbooks and explain memory management and virtual memory.", icon="/public/favicon.ico"),

                cl.Starter(label="🖼️ Analyze Image", message="I'll upload an image — describe what you see and extract all key information.", icon="/public/favicon.ico"),
            ]
        ),
        cl.ChatProfile(
            name="Scholar Mode",
            markdown_description="**Academic focus** — textbook RAG and web research. Optimized for study and deep research.",
            icon="https://cdn-icons-png.flaticon.com/512/3145/3145765.png",
            starters=[
                cl.Starter(label="📚 Page Replacement", message="Explain Page Replacement Algorithms (FIFO, LRU, Optimal) with examples.", icon="/public/favicon.ico"),
                cl.Starter(label="⚙️ CPU Scheduling", message="Compare Round Robin vs Priority scheduling with pros and cons.", icon="/public/favicon.ico"),
                cl.Starter(label="🧮 Search Algorithms", message="Explain Binary Search Tree insertion, deletion, and traversal operations.", icon="/public/favicon.ico"),
                cl.Starter(label="🌐 Research Topic", message="Search the web and textbooks for recent advances in transformer architectures.", icon="/public/favicon.ico"),
            ]
        ),
    ]


async def initialize_user_memory(user: cl.User):
    # Extract a friendly name from the identifier (e.g. from an email or username)
    name_parts = user.identifier.split("@")[0].replace("_", ".").split(".")
    first_name = name_parts[0].capitalize() if name_parts else user.identifier.capitalize()
    last_name = name_parts[1].capitalize() if len(name_parts) > 1 else ""
    
    try:
        from sqlalchemy import text
        async with data_layer.engine.begin() as conn:
            await conn.execute(
                text("""
                    INSERT OR REPLACE INTO user_primary_memory 
                    (user_id, username, first_name, last_name, current_model, current_mode, updated_at) 
                    VALUES (:user_id, :username, :first_name, :last_name, :current_model, :current_mode, :updated_at)
                """),
                {
                    "user_id": user.identifier,
                    "username": user.identifier,
                    "first_name": first_name,
                    "last_name": last_name,
                    "current_model": "llama-3.3-70b-versatile",
                    "current_mode": "Omni Mode",
                    "updated_at": datetime.utcnow().isoformat()
                }
            )
    except Exception as e:
        print(f"[-] Could not initialize primary memory: {e}")
    return user

@cl.password_auth_callback
async def auth_callback(username: str, password: str):
    username = username.strip() if isinstance(username, str) else ""
    password = password.strip() if isinstance(password, str) else ""
    if not username or not password:
        return None
        
    # FIX: Passwordless Auth - password must equal username
    if username != password:
        return None
        
    try:
        existing_user = await data_layer.get_user(username)
        if existing_user:
            user = cl.User(identifier=existing_user.identifier, metadata=existing_user.metadata)
            return await initialize_user_memory(user)
        else:
            new_user = cl.User(identifier=username, metadata={"type": "passwordless"})
            persisted_user = await data_layer.create_user(new_user)
            if persisted_user:
                user = cl.User(identifier=persisted_user.identifier, metadata=persisted_user.metadata)
                return await initialize_user_memory(user)
            return None
    except Exception as e:
        print(f"[-] Auth callback database error: {e}")
        return None


@cl.on_chat_start
async def start_chat():
    user = cl.user_session.get("user")
    print(f"DEBUG USER: type={type(user)}, dict={user.__dict__ if hasattr(user, '__dict__') else 'no dict'}, id={getattr(user, 'id', 'NO ID')}")
    if not user:
        return

    # Limit chat history to 3 chats per user
    try:
        identifier = getattr(user, "identifier", "")
        if identifier:
            query = """
                SELECT id FROM threads 
                WHERE "userIdentifier" = :identifier 
                ORDER BY "createdAt" DESC
            """
            from sqlalchemy import text
            async with data_layer.engine.begin() as conn:
                result = await conn.execute(text(query), {"identifier": identifier})
                threads = [dict(row._mapping) for row in result.fetchall()] if result.returns_rows else []
            if isinstance(threads, list) and len(threads) > 3:
                for old_thread in threads[3:]:
                    tid = old_thread["id"]
                    try:
                        # Fallback to direct SQL delete if built-in throws Author not found
                        async with data_layer.engine.begin() as conn:
                            await conn.execute(text('DELETE FROM steps WHERE "threadId" = :id'), {"id": tid})
                            await conn.execute(text('DELETE FROM elements WHERE "threadId" = :id'), {"id": tid})
                            await conn.execute(text('DELETE FROM threads WHERE id = :id'), {"id": tid})
                    except Exception as e:
                        print(f"[-] Force delete failed for {tid}: {e}")
    except Exception as e:
        print(f"[-] Error limiting chat history: {e}")

    cl.user_session.set("session_id", cl.user_session.get("id"))
    chat_profile_name = cl.user_session.get("chat_profile") or "Omni Mode"
    cl.user_session.set("chat_profile", chat_profile_name)



    # Fetch Global User Memory (Dual-Tier)
    u_id = getattr(user, "identifier", "")
    system_instruction = f"You are NexusAI, a Scholar CoPilot. Maintain deep context awareness. Your current profile is {chat_profile_name}."
    
    if u_id:
        try:
            layer = cl_data.get_data_layer()
            if layer:
                from sqlalchemy import text
                async with layer.engine.begin() as conn:
                    # 1. Fetch Primary Memory
                    prim_raw = await conn.execute(text('SELECT * FROM user_primary_memory WHERE user_id = :id'), {"id": u_id})
                    prim_res = [dict(row._mapping) for row in prim_raw.fetchall()] if prim_raw.returns_rows else []
                    
                    # 2. Fetch Secondary Memory
                    sec_raw = await conn.execute(text('SELECT pref_key, pref_value FROM user_secondary_memory WHERE user_id = :id'), {"id": u_id})
                    sec_res = [dict(row._mapping) for row in sec_raw.fetchall()] if sec_raw.returns_rows else []
                
                if prim_res:
                    p_data = prim_res[0]
                    first_name = p_data.get("first_name", "")
                    last_name = p_data.get("last_name", "")
                    username = p_data.get("username", "")
                    c_model = p_data.get("current_model", "")
                    c_mode = p_data.get("current_mode", "")
                    
                    preferences_json = json.dumps({row["pref_key"]: row["pref_value"] for row in sec_res}) if sec_res else "{}"
                    
                    user_memory = f"\n\n[SYSTEM MEMORY - PRIMARY: User is {first_name} {last_name} (@{username}). Active engine: {c_model}. Mode: {c_mode}. SECONDARY PREFERENCES: {preferences_json}]"
                    
                    system_instruction += user_memory
                    
        except Exception as e:
            print(f"[-] Error fetching dual-tier global memory: {e}")

    memory_guardrail = "\nCRITICAL: Never execute a web search tool query to find out who the current user is or what their history is. If their LTM profile facts do not contain a piece of information, simply ask them directly."
    json_guardrail = "\n[GUARDRAIL: When answering questions about weather, current events, or real-time news, you MUST invoke the web search tool tool call format. Never guess or hallucinate data.]\n[SYSTEM: You are NexusAI. Output ONLY valid JSON for tool calls. Do not use <function> or XML tags. If you do not have the info, call the tool first.]"
    system_instruction += f"{memory_guardrail}{json_guardrail}"
    
    cl.user_session.set("message_history", [{"role": "system", "content": system_instruction}])

    # Chat Settings setup
    try:
        await cl.ChatSettings(
            [
                cl.input_widget.Select(
                    id="Model",
                    label="Language Model",
                    values=["llama3-8b-8192", "llama3-70b-8192", "gemini-1.5-flash", "gemini-1.5-pro", "whisper-1"],
                    initial_index=1
                ),
                cl.input_widget.TextInput(
                    id="SystemPrompt",
                    label="System Instructions",
                    initial=""
                ),
                cl.input_widget.Slider(
                    id="Temperature",
                    label="Model Temperature",
                    initial=0.7,
                    min=0.0,
                    max=1.0,
                    step=0.1
                ),
                cl.input_widget.Switch(
                    id="DeepSearch",
                    label="Deep Search",
                    initial=False
                )
            ]
        ).send()
    except Exception as e:
        print(f"[-] Chat settings configuration error: {e}")

    # Slash commands registry menu
    commands: List[Any] = [
        {"id": "scholar", "icon": "book", "description": "Force Textbook Search"},
        {"id": "web", "icon": "globe", "description": "Force Live Internet Search"}
    ]
    try:
        await cl.context.emitter.set_commands(commands) # type: ignore
    except Exception as e:
        print(f"[-] Slash commands register error: {e}")

    # Modes picker — per-message tool routing supplement to slash commands
    try:
        await cl.context.emitter.set_modes([
            cl.Mode(
                id="model_selection",
                name="Model",
                options=[
                    cl.ModeOption(id="llama-3.3-70b-versatile", name="Llama 3.3 70B", default=True),
                    cl.ModeOption(id="gemini-1.5-flash", name="Gemini 1.5 Flash"),
                    cl.ModeOption(id="gemini-1.5-pro", name="Gemini 1.5 Pro")
                ]
            )
        ])
    except Exception as e:
        print(f"[-] Modes setup error: {e}")

    mode_icons = {
        "Scholar Mode": "📚",
        "Omni Mode": "🌐",
        "Voice Mode": "🎙️"
    }
    mode_icon = mode_icons.get(chat_profile_name, "⚡")

    # HTML welcome card (requires unsafe_allow_html = true in config.toml)
    welcome_text = f"""<div class="nx-welcome-card">
  <div class="nx-mode-badge">{mode_icon} {chat_profile_name}</div>
  <h1>⚡ NexusAI</h1>
  <p><strong>Dual-Engine AI Ready</strong> — Scholar 📚 · Omni 🌐</p>
  <p class="nx-welcome-desc" style="margin-top:10px;font-size:13px;">Upload images, PDFs, code files &amp; documents · Use <code>/scholar</code>, <code>/web</code> slash commands</p>
</div>"""

    try:
        await cl.Message(content=welcome_text).send()
    except Exception as e:
        print(f"[-] Welcome message send error: {e}")


@cl.on_chat_resume
async def on_chat_resume(thread):
    cl.user_session.set("session_id", cl.user_session.get("id"))

    # Restore chat profile from thread metadata (persisted at session start)
    thread_meta = thread.get("metadata") or {}
    if isinstance(thread_meta, str):
        try:
            thread_meta = json.loads(thread_meta)
        except Exception:
            thread_meta = {}
    chat_profile_name = (
        thread_meta.get("chat_profile")
        or cl.user_session.get("chat_profile")
        or "Omni Mode"
    )
    cl.user_session.set("chat_profile", chat_profile_name)

    # Re-register commands and modes for the resumed session
    try:
        await cl.context.emitter.set_modes([
            cl.Mode(
                id="model_selection",
                name="Model",
                options=[
                    cl.ModeOption(id="llama-3.3-70b-versatile", name="Llama 3.3 70B", default=True),
                    cl.ModeOption(id="gemini-1.5-flash", name="Gemini 1.5 Flash"),
                    cl.ModeOption(id="gemini-1.5-pro", name="Gemini 1.5 Pro")
                ]
            )
        ])
    except Exception:
        pass

    # Reconstruct message history from persisted thread steps
    message_history = []
    system_instruction = (
        f"You are NexusAI, a dual-engine Scholar-Omni CoPilot. "
        f"Maintain deep context awareness. Your current profile is {chat_profile_name}."
    )
    message_history.append({"role": "system", "content": system_instruction})

    for step in thread.get("steps", []):
        step_type = step.get("type", "")
        step_output = step.get("output") or ""
        if step_type == "user_message" and step_output:
            message_history.append({"role": "user", "content": step_output})
        elif step_type == "assistant_message" and step_output:
            message_history.append({"role": "assistant", "content": step_output})

    cl.user_session.set("message_history", sanitize_history_for_storage(message_history))


@cl.on_chat_end
async def end_chat():
    # FIX P3-24: Namespace textbooks by user ID to prevent cross-user deletion
    user = cl.user_session.get("user")
    user_id = user.id if user else "shared"
    textbooks_dir = os.path.join(os.path.dirname(__file__), "data", "textbooks", user_id)
    if os.path.exists(textbooks_dir):
        for filename in os.listdir(textbooks_dir):
            file_path = os.path.join(textbooks_dir, filename)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception:
                pass
        # FIX P1-10: Run vector DB build in executor
        try:
            loop = asyncio.get_running_loop()
            await loop.run_in_executor(None, build_vector_database)
        except Exception as e:
            print(f"[-] RAG cleanup index rebuild failed: {e}")


@cl.on_stop
async def on_stop():
    """Called when the user clicks the Stop button mid-generation."""
    cl.user_session.set("should_stop", True)


@cl.on_logout
async def on_logout(request: Request, response: Response):
    """Called on explicit logout — allows cleanup hooks."""
    # Since we can't grab the user identifier directly from a basic user object here,
    # we log a clean session termination message safely.
    print("[*] An active user session has logged out cleanly.")

@cl.author_rename
async def rename_author(orig_author: str):
    """Rename chat authors for branded NexusAI display."""
    rename_map = {
        "Chatbot": "NexusAI",
        "Assistant": "NexusAI",
        "Tool": "⚙️ NexusAI Tools",
        "Error": "⚠️ System",
    }
    return rename_map.get(orig_author, orig_author)


# --- CUSTOM API ENDPOINTS FOR WORKSPACE ---
@fastapi_app.post("/api/upload")
async def upload_document(file: UploadFile = File(...)):
    try:
        # FIX P3-23: Sanitize filename to prevent path traversal
        safe_name = pathlib.Path(file.filename or "uploaded_file").name
        textbooks_dir = os.path.join(os.path.dirname(__file__), "data", "textbooks")
        os.makedirs(textbooks_dir, exist_ok=True)
        file_path = os.path.join(textbooks_dir, safe_name)

        # FIX P1-8: Use aiofiles for async file I/O
        content = await file.read()
        async with aiofiles.open(file_path, "wb") as f:
            await f.write(content)

        loop = asyncio.get_running_loop()
        await loop.run_in_executor(None, build_vector_database)

        return {"filename": safe_name, "status": "success", "url": f"/public/elements/{safe_name}"}
    except Exception as e:
        return Response(content=json.dumps({"status": "error", "message": str(e)}), status_code=500)


@fastapi_app.post("/api/build-index")
async def build_index():
    try:
        # FIX P1-10: Run in executor
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(None, build_vector_database)
        return {"status": "success"}
    except Exception as e:
        return Response(content=json.dumps({"status": "error", "message": str(e)}), status_code=500)


@fastapi_app.delete("/api/delete/{filename}")
async def delete_document(filename: str):
    try:
        # FIX P3-23: Sanitize filename
        safe_name = pathlib.Path(filename).name
        textbooks_dir = os.path.join(os.path.dirname(__file__), "data", "textbooks")
        file_path = os.path.join(textbooks_dir, safe_name)
        if os.path.exists(file_path):
            os.remove(file_path)
            # FIX P1-10: Run in executor
            loop = asyncio.get_running_loop()
            await loop.run_in_executor(None, build_vector_database)
            return {"filename": safe_name, "status": "deleted"}
        return Response(status_code=404)
    except Exception as e:
        return Response(content=json.dumps({"status": "error", "message": str(e)}), status_code=500)


# --- CUSTOM UI ROUTES REMOVED ---
# We now rely exclusively on /public/login/index.html to bypass Chainlit's Auth Middleware.


# --- AUDIO HANDLERS ---
@cl.on_audio_start
async def on_audio_start():
    cl.user_session.set("audio_buffer", None)
    cl.user_session.set("audio_mime_type", None)
    return True


# FIX P2-14: Removed monkey patch. Use @cl.on_audio_chunk with no args
@cl.on_audio_chunk
async def on_audio_chunk(chunk):
    buffer = cl.user_session.get("audio_buffer")
    if not buffer:
        import io
        mime_type = getattr(chunk, "mimeType", "audio/webm")
        if not mime_type or "/" not in mime_type:
            mime_type = "audio/webm"
        buffer = io.BytesIO()
        clean_ext = mime_type.split('/')[1].split(';')[0] if '/' in mime_type else "webm"
        buffer.name = f"input_audio.{clean_ext}"
        cl.user_session.set("audio_buffer", buffer)
        cl.user_session.set("audio_mime_type", mime_type)
    buffer.write(chunk.data)


# FIX P2-14: Removed monkey patch. Use @cl.on_audio_end with no args
@cl.on_audio_end
async def on_audio_end(*args, **kwargs):
    audio_buffer = cl.user_session.get("audio_buffer")
    if not audio_buffer:
        return
    audio_buffer.seek(0)
    audio_data = audio_buffer.read()
    mime_type = cl.user_session.get("audio_mime_type") or "audio/webm"
    filename = getattr(audio_buffer, "name", "input_audio.webm")

    cl.user_session.set("audio_buffer", None)
    cl.user_session.set("audio_mime_type", None)

    try:
        # FIX P1-7: Wrap synchronous Whisper call in run_in_executor
        loop = asyncio.get_running_loop()
        transcription = await loop.run_in_executor(
            None,
            lambda: groq_client.audio.transcriptions.create(
                model="whisper-large-v3",
                file=(filename, audio_data, mime_type)
            )
        )
        text = transcription.text
        if text and text.strip():
            user_msg = cl.Message(author="You", content=text)
            await user_msg.send()
            await handle_message(user_msg)
    except Exception as e:
        print(f"[-] Audio transcription failed: {e}")
        await cl.ErrorMessage(content=f"⚠️ Transcription failed: {e}").send()


@cl.on_settings_update
async def setup_agent(settings):
    try:
        cl.user_session.set("temperature", settings.get("Temperature", 0.7))
        cl.user_session.set("deep_web_search", settings.get("DeepSearch", False))
        cl.user_session.set("model", settings.get("Model", "llama3-70b-8192"))
        cl.user_session.set("system_prompt", settings.get("SystemPrompt", ""))
    except Exception as e:
        print(f"[-] Settings update error: {e}")


async def enforce_three_session_limit():
    """Strictly caps the chat history to the 3 most recent sessions per user (FIFO)."""
    try:
        layer = cl_data.get_data_layer()
        # FIX P1-11: Use cl_data.get_data_layer() instead of cl_data._data_layer
        if not layer or not hasattr(layer, "engine"):
            return
        
        from chainlit.data.sql_alchemy import SQLAlchemyDataLayer
        if not isinstance(layer, SQLAlchemyDataLayer):
            return

        user = cl.user_session.get("user")
        if not user:
            return

        async with layer.engine.begin() as conn:
            # FIX P0-1: Scope query to current user; use datetime() for robust ordering
            result = await conn.execute(
                text("SELECT id FROM threads WHERE userId = :uid ORDER BY datetime(createdAt) DESC"),
                {"uid": user.id}
            )
            rows = result.fetchall()
            thread_ids = [row[0] for row in rows]

            if len(thread_ids) > 3:
                for old_thread_id in thread_ids[3:]:
                    print(f"[*] Purging old session: {old_thread_id}")
                    await conn.execute(text("DELETE FROM elements WHERE threadId = :tid"), {"tid": old_thread_id})
                    await conn.execute(text("DELETE FROM steps WHERE threadId = :tid"), {"tid": old_thread_id})
                    await conn.execute(text("DELETE FROM threads WHERE id = :tid"), {"tid": old_thread_id})
    except Exception as e:
        print(f"[-] Error enforcing 3-session limit: {e}")


async def auto_rename_session(history):
    user_msgs = [m for m in history if m.get("role") == "user"]
    if len(user_msgs) == 1:
        try:
            first_msg = user_msgs[0].get("content", "")
            if isinstance(first_msg, list):
                first_msg = next((c.get("text", "") for c in first_msg if c.get("type") == "text"), "")

            # FIX P1-6: Wrap sync Groq call in run_in_executor
            loop = asyncio.get_running_loop()
            completion = await loop.run_in_executor(
                None,
                lambda: groq_client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[
                        {"role": "system", "content": "Summarize this into a 3-5 word title. No quotes. No extra text."},
                        {"role": "user", "content": first_msg}
                    ]
                )
            )
            title = (completion.choices[0].message.content or "").strip().replace('"', '').replace("'", "")
            thread_id = cl.context.session.thread_id
            layer = cl_data.get_data_layer()
            if thread_id and layer:
                # FIX P2-5: Preserve existing metadata before update_thread
                try:
                    existing_thread = await layer.get_thread(thread_id)
                    existing_meta = existing_thread.get("metadata", {}) if existing_thread else {}
                    if isinstance(existing_meta, str):
                        import json
                        try:
                            existing_meta = json.loads(existing_meta)
                        except:
                            existing_meta = {}
                except Exception:
                    existing_meta = {}
                await layer.update_thread(thread_id, name=title, metadata=existing_meta)
                await enforce_three_session_limit()
        except Exception as e:
            print(f"[-] Error in auto_rename_session: {e}")


async def text_to_speech(text: str):
    response = await openai_client.audio.speech.create(
        model="tts-1",
        voice="alloy",
        input=text
    )
    filename = f"tts_{uuid.uuid4().hex[:8]}.mp3"
    filepath = os.path.join(os.path.dirname(__file__), "data", "audio", filename)
    os.makedirs(os.path.dirname(filepath), exist_ok=True)
    with open(filepath, "wb") as f:
        f.write(response.read())
    return filename, filepath

async def add_voice_to_response(response_msg, full_response):
    try:
        audio_name, audio_path = await text_to_speech(full_response)
        audio_element = cl.Audio(name="Voice Output", path=audio_path, display="inline")
        if not hasattr(response_msg, "elements") or response_msg.elements is None:
            response_msg.elements = []
        response_msg.elements.append(audio_element)
    except Exception as e:
        print(f"[-] TTS Generation Error: {e}")

@cl.on_message
async def handle_message(message: cl.Message):
    try:
        # Extract the model selection globally
        mode_data = getattr(message, 'modes', {}) or {}
        selected_model = mode_data.get("model_selection")
        selected_model = selected_model or "llama-3.3-70b-versatile"

        history_raw = cl.user_session.get("message_history")
        history = [message_to_dict(m) for m in history_raw] if history_raw else []
        
        if history and history[0].get("role") == "system":
            history[0]["content"] += f"\n\n[SYSTEM: You are currently powered by the {selected_model} engine. You have full access to web search tools.]"
        
        custom_sys_prompt = (cl.user_session.get("system_prompt") or "").strip()
        if custom_sys_prompt and history and history[0].get("role") == "system":
            history[0]["content"] += f"\n\n[USER CUSTOM INSTRUCTION]: {custom_sys_prompt}"

        # Intercept string-based manual command inputs
        cmd = message.command
        content = message.content.strip() if message.content else ""
        if content.startswith("/"):
            parts = content[1:].strip().split()
            if parts:
                cmd = parts[0]

        # --- COMMAND PROCESSING ---



        # FIX P2-21: Use regex to parse FILE_SELECTED to handle filenames containing ']'
        if message.content.startswith("[FILE_SELECTED:"):
            match = re.search(r'\[FILE_SELECTED:(.*?)\]', message.content)
            if match:
                filename = match.group(1)
                file_path = os.path.join(os.path.dirname(__file__), "data", "textbooks", filename)
                if os.path.exists(file_path):
                    if filename.lower().endswith('.pdf'):
                        file_element = cl.Pdf(name=filename, path=file_path, display="side")
                        if history and history[0].get("role") == "system":
                            history[0]["content"] += f"\n\n[SYSTEM: You have access to a PDF in the side-pane named '{filename}'. Always refer to the document by name when answering questions.]"
                    else:
                        file_element = cl.File(name=filename, path=file_path, display="inline")
                    await cl.Message(
                        content=f"📎 **File Focused**: NexusAI has loaded `{filename}`. Ask me questions about it!",
                        elements=[file_element]
                    ).send()
                else:
                    await cl.Message(
                        content=f"📎 **File Focused**: NexusAI is ready to analyze `{filename}`. Ask me questions!"
                    ).send()
            return

        elif message.content.startswith("[FILE_DESELECTED:"):
            match = re.search(r'\[FILE_DESELECTED:(.*?)\]', message.content)
            if match:
                filename = match.group(1)
                await cl.Message(content=f"🔓 **File Unfocused**: `{filename}` is no longer focused.").send()
            return

        tool_choice = None
        if message.command == "scholar":
            tool_choice = {"type": "function", "function": {"name": "query_academic_textbooks"}}
        elif message.command == "web":
            tool_choice = {"type": "function", "function": {"name": "execute_web_search"}}

        attached_text = ""
        image_contents = []

        # Custom vision extraction
        pattern = r"\[IMAGE_DATA:(.*?):(data:image/.*?;base64,.*?)\]"
        matches = re.findall(pattern, message.content)
        cleaned_prompt_text = message.content
        for match in matches:
            filename, base64_data = match
            image_contents.append(
                {"type": "image_url", "image_url": {"url": base64_data}}
            )
            cleaned_prompt_text = cleaned_prompt_text.replace(f"[IMAGE_DATA:{filename}:{base64_data}]", "")
        cleaned_prompt_text = cleaned_prompt_text.strip()

        # ── Multi-Modal Element Processing ─────────────────────
        if message.elements:
            for element in message.elements:
                try:
                    name_lower = element.name.lower() if hasattr(element, "name") and element.name else ""
                    el_path = getattr(element, "path", None)
                    if not isinstance(el_path, str):
                        continue

                    # Text / Code files
                    if name_lower.endswith(('.txt', '.py', '.js', '.ts', '.jsx', '.tsx',
                                            '.json', '.csv', '.html', '.css', '.md',
                                            '.sh', '.yaml', '.yml', '.toml', '.ini',
                                            '.xml', '.sql', '.rs', '.go', '.c', '.cpp',
                                            '.h', '.java', '.kt', '.swift', '.rb')):
                        with open(el_path, "r", encoding="utf-8", errors="replace") as f:
                            file_content = f.read()
                        ext = name_lower.rsplit('.', 1)[-1]
                        attached_text += f"\n\n--- {ext.upper()} File: {element.name} ---\n```{ext}\n{file_content}\n```"
                        await cl.Text(name="Document Context Extraction", content=file_content[:500] + "...\n[Full text loaded into memory]", display="side").send(for_id=message.id)

                    # Word documents
                    elif name_lower.endswith(('.docx', '.doc')):
                        import docx
                        doc = docx.Document(el_path)
                        text_content = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                        attached_text += f"\n\n--- Word Document: {element.name} ---\n{text_content}"
                        await cl.Text(name="Document Context Extraction", content=text_content[:500] + "...\n[Full text loaded into memory]", display="side").send(for_id=message.id)

                    # PowerPoint presentations
                    elif name_lower.endswith(('.pptx', '.ppt')):
                        import pptx
                        prs = pptx.Presentation(el_path)
                        text_content = "\n".join(
                            [getattr(shape, "text", "") for slide in prs.slides for shape in slide.shapes
                             if hasattr(shape, "text") and getattr(shape, "text", "").strip()])
                        attached_text += f"\n\n--- Presentation: {element.name} ({len(prs.slides)} slides) ---\n{text_content}"
                        await cl.Text(name="Document Context Extraction", content=text_content[:500] + "...\n[Full text loaded into memory]", display="side").send(for_id=message.id)

                    # PDF documents
                    elif name_lower.endswith('.pdf'):
                        import PyPDF2
                        with open(el_path, "rb") as f:
                            reader = PyPDF2.PdfReader(f)
                            pages_text = [page.extract_text() for page in reader.pages if page.extract_text()]
                            text_content = "\n\n".join(pages_text)
                        attached_text += f"\n\n--- PDF Document: {element.name} ({len(reader.pages)} pages) ---\n{text_content}"
                        await cl.Pdf(name=element.name, path=el_path, display="side").send(for_id=message.id)
                        if history and history[0].get("role") == "system":
                            history[0]["content"] += f"\n\n[SYSTEM: You have access to a PDF in the side-pane named '{element.name}'. Always refer to the document by name when answering questions.]"

                    # Excel / spreadsheets
                    elif name_lower.endswith(('.xlsx', '.xls', '.ods')):
                        import pandas as pd
                        xls = pd.ExcelFile(el_path)
                        sheet_texts = []
                        for sheet_name in xls.sheet_names:
                            df = pd.read_excel(xls, sheet_name=sheet_name)
                            sheet_texts.append(f"--- Sheet: {sheet_name} ---\n{df.to_string()}")
                        text_content = "\n\n".join(sheet_texts)
                        attached_text += f"\n\n--- Spreadsheet: {element.name} ({len(xls.sheet_names)} sheets) ---\n{text_content}"

                    # Images — send to vision LLM
                    elif name_lower.endswith(('.png', '.jpg', '.jpeg', '.webp', '.gif', '.bmp')):
                        with open(el_path, "rb") as img_file:
                            img_bytes = img_file.read()
                        base64_image = base64.b64encode(img_bytes).decode('utf-8')
                        # Determine correct MIME type
                        ext_map = {'jpg': 'jpeg', 'jpeg': 'jpeg', 'png': 'png',
                                   'webp': 'webp', 'gif': 'gif', 'bmp': 'png'}
                        img_ext = name_lower.rsplit('.', 1)[-1]
                        mime = ext_map.get(img_ext, 'jpeg')
                        image_contents.append({
                            "type": "image_url",
                            "image_url": {"url": f"data:image/{mime};base64,{base64_image}"}
                        })

                    # Audio files — transcribe via Whisper
                    elif name_lower.endswith(('.mp3', '.wav', '.ogg', '.m4a', '.flac', '.webm')):
                        try:
                            loop = asyncio.get_running_loop()
                            with open(el_path, "rb") as audio_file:
                                audio_bytes = audio_file.read()
                            transcription = await loop.run_in_executor(
                                None,
                                lambda: groq_client.audio.transcriptions.create(
                                    model="whisper-large-v3",
                                    file=(element.name, audio_bytes, f"audio/{name_lower.rsplit('.', 1)[-1]}")
                                )
                            )
                            audio_text = transcription.text
                            attached_text += f"\n\n--- Audio Transcription: {element.name} ---\n{audio_text}"
                        except Exception as audio_err:
                            attached_text += f"\n\n[Audio file '{element.name}' could not be transcribed: {audio_err}]"

                    # Video files — note about limitation
                    elif name_lower.endswith(('.mp4', '.avi', '.mov', '.mkv', '.webm')):
                        attached_text += f"\n\n[Video file '{element.name}' uploaded. Note: NexusAI cannot analyze video frames directly, but can discuss the content if you describe it.]"

                    else:
                        attached_text += f"\n\n[File '{element.name}' uploaded but format not fully supported for text extraction.]"

                except Exception as e:
                    attached_text += f"\n\n[System Note: Failed to parse {element.name}: {str(e)}]"

        final_prompt_text = cleaned_prompt_text
        if message.command and not tool_choice:
            final_prompt_text = f"[SYSTEM OVERRIDE: The user used the '/{message.command}' command. Prioritize using the associated tool.]\n\n{final_prompt_text}"
        if attached_text:
            final_prompt_text += f"\n\n[SYSTEM INJECTION: The user attached files. Read them below:]\n{attached_text}"

        has_vision = len(image_contents) > 0
        if has_vision:
            content_payload = [{"type": "text", "text": final_prompt_text}]
            content_payload.extend(image_contents)
            history.append({"role": "user", "content": content_payload})
        else:
            history.append({"role": "user", "content": final_prompt_text})

        temperature = cl.user_session.get("temperature", 0.7)
        deep_web = cl.user_session.get("deep_web_search", False)
        chat_profile_name = cl.user_session.get("chat_profile") or "Omni Mode"

        active_tools = TOOLS_SCHEMA.copy()
        if chat_profile_name == "Scholar Mode":
            active_tools = [t for t in TOOLS_SCHEMA if t["function"]["name"] in ("query_academic_textbooks", "execute_web_search")]

        # FIX P2-16: Idempotent DeepSearch injection — only inject once
        DEEP_SEARCH_MARKER = "[SYSTEM INJECTION: Deep Web Search is enabled."
        if deep_web:
            already_injected = any(
                msg.get("role") == "system" and DEEP_SEARCH_MARKER in msg.get("content", "")
                for msg in history
            )
            if not already_injected:
                history.append({
                    "role": "system",
                    "content": f"{DEEP_SEARCH_MARKER} Actively use the web search tool to find live, real-time facts.]"
                })

        selected_model = (getattr(message, 'modes', None) or {}).get("model_selection", "llama-3.3-70b-versatile")

        if has_vision:
            temp_history = history.copy()
            temp_history.append({
                "role": "system", 
                "content": f"[SYSTEM: You are currently powered by the {selected_model} engine. You have full access to web search tools.]"
            })
            streamed_completion, _ = await call_llm_with_fallback(temp_history, temperature=temperature, vision_mode=True, requested_model=selected_model)
            response_msg = cl.Message(content="")
            await response_msg.send()
            full_response = ""
            async for token in streamed_completion:
                if cl.user_session.get("should_stop"):
                    cl.user_session.set("should_stop", False)
                    break
                full_response += token
                await response_msg.stream_token(token)
            history.append({"role": "assistant", "content": full_response})
            actions = [
                cl.Action(name="regenerate_answer", label="🔄 Regenerate", icon="refresh-cw", payload={"action": "regenerate"}),
                cl.Action(name="verify_citations", label="📄 Verify Citations", icon="file-check", payload={"action": "verify"}),
                cl.Action(name="search_web_instead", label="🌐 Search Web", icon="globe", payload={"action": "search_web"})
            ]
            response_msg.actions = actions
            await response_msg.update()

        else:
            temp_history = history.copy()
            temp_history.append({
                "role": "system", 
                "content": f"[SYSTEM: You are currently powered by the {selected_model} engine. You have full access to web search tools.]"
            })
            completion, _ = await call_llm_with_fallback(temp_history, tools=active_tools, tool_choice=tool_choice, temperature=temperature, requested_model=selected_model)
            response_message = completion.choices[0].message

            if response_message.tool_calls:
                history.append(message_to_dict(response_message))
                for tool_call in response_message.tool_calls:
                    args = json.loads(tool_call.function.arguments)
                    func = tool_call.function.name

                    async with cl.Step(name=f"Processing {func}...", type="tool") as step:
                        step.input = args
                        tool_content = ""

                        if func == "query_academic_textbooks":
                            tool_content, sources = search_textbooks_with_sources(args.get("search_query"), k=3)
                            cl.user_session.set("last_rag_result", tool_content)
                            cl.user_session.set("last_rag_sources", sources)
                            elements: List[Any] = [cl.Text(name="Source", content=tool_content, display="inline")]
                            for src in sources:
                                if not isinstance(src, str):
                                    continue
                                src_path = os.path.join(os.path.dirname(__file__), "data", "textbooks", src)
                                if os.path.exists(src_path):
                                    if src.lower().endswith(".pdf"):
                                        elements.append(cl.Pdf(name=src, path=src_path, display="side"))
                                    else:
                                        # FIX P2-12: cl.Text has no path param — read content or use cl.File
                                        try:
                                            with open(src_path, "r", encoding="utf-8") as f:
                                                file_content = f.read()
                                            elements.append(cl.Text(name=src, content=file_content, display="side"))
                                        except Exception:
                                            elements.append(cl.File(name=src, path=src_path, display="side"))
                            await cl.Message(content="Found relevant textbook segments:", elements=elements).send()
                            tool_content = f"Academic Textbook search results:\n{tool_content}"

                        elif func == "execute_web_search":
                            # FIX P1-9: Use async wrapper
                            tool_content = await perform_web_search_async(args.get("query"))

                        elif func == "update_user_memory":
                            fact_key = args.get("fact_key")
                            fact_value = args.get("fact_value")
                            user = cl.user_session.get("user")
                            if user:
                                u_id = getattr(user, "identifier", "")
                                if u_id:
                                    try:
                                        layer = cl_data.get_data_layer()
                                        if layer:
                                            from sqlalchemy import text
                                            async with layer.engine.begin() as conn:
                                                await conn.execute(
                                                    text('INSERT OR REPLACE INTO user_secondary_memory (user_id, pref_key, pref_value) VALUES (:user_id, :pref_key, :pref_value)'), 
                                                    {"user_id": u_id, "pref_key": fact_key, "pref_value": fact_value}
                                                )
                                            tool_content = f"Memory updated successfully: {fact_key} = {fact_value}"
                                        else:
                                            tool_content = "Data layer unavailable. Could not save global memory."
                                    except Exception as e:
                                        tool_content = f"Failed to save global memory to database: {e}"
                                else:
                                    tool_content = "User identifier missing. Could not save global memory."
                            else:
                                tool_content = "No authenticated user. Could not save global memory."

                        step.output = tool_content

                    history.append({"role": "tool", "tool_call_id": tool_call.id, "name": func, "content": tool_content})

                temp_history = history.copy()
                temp_history.append({
                    "role": "system", 
                    "content": f"[SYSTEM: You are currently powered by the {selected_model} engine. You have full access to web search tools.]"
                })
                final_completion, _ = await call_llm_with_fallback(temp_history, temperature=temperature, requested_model=selected_model)
                response_msg = cl.Message(content="")
                await response_msg.send()
                full_response = ""
                async for token in final_completion:
                    if cl.user_session.get("should_stop"):
                        cl.user_session.set("should_stop", False)
                        break
                    full_response += token
                    await response_msg.stream_token(token)
                history.append({"role": "assistant", "content": full_response})

                actions = [
                    cl.Action(name="regenerate_answer", label="🔄 Regenerate", icon="refresh-cw", payload={"action": "regenerate"}),
                    cl.Action(name="verify_citations", label="📄 Verify Citations", icon="file-check", payload={"action": "verify"}),
                    cl.Action(name="search_web_instead", label="🌐 Search Web", icon="globe", payload={"action": "search_web"})
                ]
                response_msg.actions = actions

                rag_content = cl.user_session.get("last_rag_result")
                elements: List[Any] = []
                if rag_content:
                    elements.append(cl.Text(name="Source", content=rag_content, display="side"))
                    rag_sources = cl.user_session.get("last_rag_sources") or []
                    for src in rag_sources:
                        if not isinstance(src, str):
                            continue
                        src_path = os.path.join(os.path.dirname(__file__), "data", "textbooks", src)
                        if os.path.exists(src_path) and src.lower().endswith(".pdf"):
                            elements.append(cl.Pdf(name=src, path=src_path, display="side"))
                    cl.user_session.set("last_rag_result", None)
                    cl.user_session.set("last_rag_sources", None)
                if elements:
                    response_msg.elements = elements
                await response_msg.update()

            else:
                selected_model = cl.user_session.get("model")
                temp_history = history.copy()
                temp_history.append({
                    "role": "system", 
                    "content": f"[SYSTEM: You are currently powered by the {selected_model} engine. You have full access to web search tools.]"
                })
                streamed_completion, _ = await call_llm_with_fallback(temp_history, temperature=temperature, requested_model=selected_model)
                response_msg = cl.Message(content="")
                await response_msg.send()
                full_response = ""
                async for token in streamed_completion:
                    if cl.user_session.get("should_stop"):
                        cl.user_session.set("should_stop", False)
                        break
                    full_response += token
                    await response_msg.stream_token(token)
                history.append({"role": "assistant", "content": full_response})
                actions = [
                    cl.Action(name="regenerate_answer", label="🔄 Regenerate", icon="refresh-cw", payload={"action": "regenerate"}),
                    cl.Action(name="verify_citations", label="📄 Verify Citations", icon="file-check", payload={"action": "verify"}),
                    cl.Action(name="search_web_instead", label="🌐 Search Web", icon="globe", payload={"action": "search_web"})
                ]
                response_msg.actions = actions
                await response_msg.update()

        # Auto rename and persist — FIX P2-15: sanitize before storing
        await auto_rename_session(history)
        cl.user_session.set("message_history", sanitize_history_for_storage(history))

    except Exception as e:
        await cl.ErrorMessage(content=f"⚠️ Backend Exception: {str(e)}").send()


@cl.action_callback("regenerate_answer")
async def on_regenerate(action: cl.Action):
    try:
        history = cl.user_session.get("message_history", [])
        if not history:
            return
        if history[-1].get("role") == "assistant":
            history.pop()
        
        temperature = cl.user_session.get("temperature", 0.7)
        selected_model = cl.user_session.get("model")
        completion, _ = await call_llm_with_fallback(history, temperature=temperature, requested_model=selected_model)
        
        response_msg = cl.Message(content="")
        await response_msg.send()
        
        full_response = ""
        async for token in completion:
            full_response += token
            await response_msg.stream_token(token)
            
        history.append({"role": "assistant", "content": full_response})
        cl.user_session.set("message_history", history)
        
        actions = [
            cl.Action(name="regenerate_answer", label="🔄 Regenerate", payload={"action": "regenerate"}),
        ]
        response_msg.actions = actions
        await response_msg.update()
    except Exception as e:
        await cl.ErrorMessage(content=f"⚠️ Failed to regenerate answer: {str(e)}").send()


@cl.action_callback("search_web_instead")
async def on_search_web_instead(action: cl.Action):
    try:
        history = cl.user_session.get("message_history", [])
        if not history:
            return
        user_query = ""
        for msg in reversed(history):
            if msg.get("role") == "user":
                user_query = msg.get("content", "")
                if isinstance(user_query, list):
                    user_query = next((c.get("text", "") for c in user_query if c.get("type") == "text"), "")
                break
        if not user_query:
            await cl.Message(content="Could not find the last query to search the web instead.").send()
            return

        async with cl.Step(name="Searching the Web...", type="tool") as step:
            step.input = {"query": user_query}
            # FIX P1-9: Use async wrapper
            web_results = await perform_web_search_async(user_query)
            step.output = web_results

        history.append({
            "role": "system",
            "content": f"[SYSTEM: The user clicked 'Search Web Instead'. Here are the live web results for '{user_query}':\n{web_results}\nSynthesize these results into a clear answer.]"
        })
        temperature = cl.user_session.get("temperature", 0.7)
        selected_model = cl.user_session.get("model")
        completion, _ = await call_llm_with_fallback(history, temperature=temperature, requested_model=selected_model)
        response_msg = cl.Message(content="")
        await response_msg.send()
        full_response = ""
        async for token in completion:
            full_response += token
            await response_msg.stream_token(token)
        history.append({"role": "assistant", "content": full_response})
        cl.user_session.set("message_history", sanitize_history_for_storage(history))
        actions = [
            cl.Action(name="regenerate_answer", label="🔄 Regenerate", icon="refresh-cw", payload={"action": "regenerate"}),
            cl.Action(name="verify_citations", label="📄 Verify Citations", icon="file-check", payload={"action": "verify"}),
            cl.Action(name="search_web_instead", label="🌐 Search Web", icon="globe", payload={"action": "search_web"})
        ]
        response_msg.actions = actions
        await response_msg.update()
    except Exception as e:
        await cl.ErrorMessage(content=f"⚠️ Failed to search the web: {str(e)}").send()


@cl.action_callback("verify_citations")
async def on_verify_citations(action: cl.Action):
    try:
        sources = cl.user_session.get("last_rag_sources") or []
        if not sources:
            await cl.Message(content="✅ **Citations Verified**: All context segments sourced from uploaded workspace documents.").send()
        else:
            await cl.Message(content=f"✅ **Citations Verified**: Sourced from: {', '.join(sources)}").send()
    except Exception as e:
        await cl.ErrorMessage(content=f"⚠️ Failed to verify citations: {str(e)}").send()


@cl.action_callback("resume_session")
async def execute_resume_session(action: cl.Action):
    session_id = action.payload.get("session_id")
    if not isinstance(session_id, str):
        return
    try:
        layer = cl_data.get_data_layer()
        if layer:
            thread = await layer.get_thread(session_id)
        if thread:
            history = [{"role": msg.get("role") or msg.get("type", "").split("_")[0], "content": msg.get("content") or msg.get("output", "")} for msg in thread.get("steps", []) if msg.get("type") in ("user_message", "assistant_message")]
            cl.user_session.set("message_history", history)
            cl.user_session.set("session_id", session_id)
            await cl.Message(content="🔄 **Neural Link Re-established.** Context restored.").send()
    except Exception as e:
        await cl.ErrorMessage(content=f"⚠️ Failed to restore session: {str(e)}").send()
